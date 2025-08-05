import os
import re
from typing import List, Sequence, Pattern, Dict
from llama_index.core import Document
from llama_index.core.schema import BaseNode
from llama_parse import LlamaParse
from dotenv import load_dotenv

load_dotenv()

class DocumentEmbeddingMethod:
    def __init__(self, docs_path: str):
        self.docs_path = docs_path
        self.supported_extensions = ['.pdf', '.docx', '.pptx', '.html', '.txt']
        self.parser = LlamaParse(
            api_key=os.getenv("LLAMA_CLOUD_API_KEY"),
            result_type="markdown",  # Daha iyi yapılandırılmış çıktı
            num_workers=4,           # Paralel işleme
            verbose=True
        )

    @staticmethod
    def customize_metadata(document: Document, data_source_id: str) -> Document:
        document.metadata.update({
            "data_source_id": data_source_id,
            "file_type": document.metadata.get("file_name", "").split(".")[-1].lower(),
            "processing_method": "llama_parse"
        })
        return document

    def _compile_patterns(self, patterns: List[str]) -> List[Pattern]:
        compiled = []
        for pattern in patterns:
            try:
                compiled.append(re.compile(pattern))
            except re.error:
                print(f"Geçersiz regex deseni: {pattern}")
        return compiled

    def apply_rules():

        for doc in documents:
            file_path = doc.metadata.get("file_path", "")
            file_ext = doc.metadata.get("file_extension", "").lower()

            if file_ext not in [ext[1:] for ext in self.supported_extensions]:
                continue

            excluded = any(pattern.search(file_path) for pattern in compiled_exclude)
            included = any(pattern.search(file_path) for pattern in compiled_include) if compiled_include else True

            if not excluded and included:
                filtered_docs.append(doc)

        print(f"[apply_rules] Filtrelenen doküman sayısı: {len(filtered_docs)}")
        return filtered_docs

    def _process_with_llamaparse(self, file_path: str) -> str:
        """LlamaParse ile dosya işleme"""
        try:
            documents = self.parser.load_data(file_path)
            return "\n".join([doc.text for doc in documents])
        except Exception as e:
            print(f"LlamaParse işleme hatası ({file_path}): {str(e)}")
            return ""

    def _process_docx(self, file_path: str) -> str:
        """DOCX işleme (LlamaParse desteklemiyorsa)"""
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            print(f"DOCX işleme hatası ({file_path}): {str(e)}")
            return ""

    def _process_pptx(self, file_path: str) -> str:
        """PPTX işleme (LlamaParse desteklemiyorsa)"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            print(f"PPTX işleme hatası ({file_path}): {str(e)}")
            return ""

    def _process_html(self, file_path: str) -> str:
        """HTML işleme"""
        try:
            from unstructured.partition.html import partition_html
            with open(file_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            elements = partition_html(text=html_content)
            return "\n".join([str(el) for el in elements])
        except Exception as e:
            print(f"HTML işleme hatası ({file_path}): {str(e)}")
            return ""

    def _process_txt(self, file_path: str) -> str:
        """TXT işleme"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            print(f"TXT işleme hatası ({file_path}): {str(e)}")
            return ""

    def get_documents(self, data_source_id: str) -> List[Document]:
        documents = []
        print(f"\n[get_documents] Dokümanlar taranıyor: {self.docs_path}")

        for root, _, files in os.walk(self.docs_path):
            for file_name in files:
                file_ext = os.path.splitext(file_name)[1].lower()
                if file_ext not in self.supported_extensions:
                    continue

                file_path = os.path.join(root, file_name)
                content = ""
                
                try:
                    if file_ext == '.pdf':
                        content = self._process_with_llamaparse(file_path)
                    elif file_ext == '.docx':
                        content = self._process_docx(file_path)
                    elif file_ext == '.pptx':
                        content = self._process_pptx(file_path)
                    elif file_ext == '.html':
                        content = self._process_html(file_path)
                    elif file_ext == '.txt':
                        content = self._process_txt(file_path)
                except Exception as e:
                    print(f"Doküman işleme hatası ({file_path}): {str(e)}")
                    continue

                if not content.strip():
                    print(f"Uyarı: {file_path} boş içerik")
                    continue

                doc = Document(
                    text=content,
                    metadata={
                        "file_path": file_path,
                        "file_name": file_name,
                        "file_extension": file_ext[1:],
                        "last_modified": os.path.getmtime(file_path),
                        "file_size": os.path.getsize(file_path)
                    }
                )
                self.customize_metadata(doc, data_source_id)
                documents.append(doc)

        print(f"[get_documents] İşlenen doküman sayısı: {len(documents)}")
        return documents

    def get_nodes(self, documents: Sequence[Document]) -> Sequence[BaseNode]:
        return []